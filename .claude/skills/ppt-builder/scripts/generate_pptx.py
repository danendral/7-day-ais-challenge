import json
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

BLUE_DARK = RGBColor(0x1B, 0x3A, 0x5C)
BLUE_MED = RGBColor(0x2E, 0x75, 0xB6)
BLUE_LIGHT = RGBColor(0x9D, 0xC3, 0xE6)
BLUE_BG = RGBColor(0xF2, 0xF7, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GREEN = RGBColor(0x2D, 0x8C, 0x46)
RED = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x99, 0x99, 0x99)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_bg(slide, color=BLUE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=16, color=BLACK, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_accent_bar(slide, left, top, width, height, color=BLUE_MED):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_metric_card(slide, left, top, width, height, label, value, delta=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLUE_LIGHT
    shape.line.width = Pt(1)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = GRAY
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = BLUE_DARK
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    if delta:
        p3 = tf.add_paragraph()
        arrow = "▲" if delta.get("direction") == "up" else "▼"
        delta_color = GREEN if delta.get("direction") == "up" else RED
        p3.text = f"{arrow} {delta.get('value', '')}"
        p3.font.size = Pt(14)
        p3.font.bold = True
        p3.font.color.rgb = delta_color
        p3.font.name = "Calibri"
        p3.alignment = PP_ALIGN.CENTER

    return shape


def build_title_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)
    add_accent_bar(slide, 0, 3.2, 13.333, 0.08, BLUE_LIGHT)
    add_textbox(slide, 0.8, 1.5, 11.7, 1.5, slide_data.get("title", ""),
                font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if slide_data.get("subtitle"):
        add_textbox(slide, 0.8, 3.5, 11.7, 1.0, slide_data["subtitle"],
                    font_size=20, color=BLUE_LIGHT, alignment=PP_ALIGN.CENTER)
    if slide_data.get("date"):
        add_textbox(slide, 0.8, 5.5, 11.7, 0.5, slide_data["date"],
                    font_size=14, color=BLUE_LIGHT, alignment=PP_ALIGN.CENTER)


def build_content_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, BLUE_MED)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.8, slide_data.get("title", ""),
                font_size=28, bold=True, color=BLUE_DARK)
    add_accent_bar(slide, 0.8, 1.15, 2.0, 0.04, BLUE_MED)

    y = 1.5
    if slide_data.get("metrics"):
        metrics = slide_data["metrics"]
        count = len(metrics)
        card_w = min(2.8, 11.7 / max(count, 1))
        gap = 0.3
        total_w = count * card_w + (count - 1) * gap
        start_x = (13.333 - total_w) / 2
        for i, m in enumerate(metrics):
            x = start_x + i * (card_w + gap)
            add_metric_card(slide, x, y, card_w, 1.6,
                            m.get("label", ""), m.get("value", ""),
                            m.get("delta"))
        y += 2.0

    if slide_data.get("bullets"):
        add_bullet_textbox(slide, 0.8, y, 5.5, 4.5 - (y - 1.5),
                           slide_data["bullets"], font_size=15, color=BLACK)

    if slide_data.get("right_notes"):
        add_textbox(slide, 7.0, y, 5.5, 4.5 - (y - 1.5),
                    slide_data["right_notes"], font_size=15, color=GRAY)


CHART_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
}


def add_chart(slide, chart_type, chart_data_obj, left, top, width, height, title=""):
    chart_shape = slide.shapes.add_chart(
        chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data_obj
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    chart.legend.font.color.rgb = GRAY

    if hasattr(chart, 'value_axis') and chart.value_axis:
        chart.value_axis.has_title = False
        chart.value_axis.major_gridlines.format.line.color.rgb = BLUE_LIGHT
        chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.color.rgb = GRAY

    if hasattr(chart, 'category_axis') and chart.category_axis:
        chart.category_axis.tick_labels.font.size = Pt(9)
        chart.category_axis.tick_labels.font.color.rgb = GRAY

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.color.rgb = BLUE_DARK

    plot = chart.plots[0]
    plot.gap_width = 80
    series_colors = [BLUE_MED, BLUE_DARK, BLUE_LIGHT, RGBColor(0x4A, 0x90, 0xD9), RGBColor(0x6B, 0xA5, 0xE7)]
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = series_colors[i % len(series_colors)]

    return chart_shape


def build_chart_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_BG)
    add_accent_bar(slide, 0, 0, 13.333, 0.06, BLUE_MED)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.8, slide_data.get("title", ""),
                font_size=28, bold=True, color=BLUE_DARK)
    add_accent_bar(slide, 0.8, 1.15, 2.0, 0.04, BLUE_MED)

    charts = slide_data.get("charts", [])
    if not charts:
        return

    if len(charts) == 1:
        cd = charts[0]
        chart_type_key = cd.get("type", "bar")
        xl_type = CHART_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)
        data = CategoryChartData()
        data.categories = cd.get("categories", [])
        for s in cd.get("series", []):
            data.add_series(s.get("name", ""), s.get("values", []))
        add_chart(slide, xl_type, data, 0.8, 1.5, 11.7, 5.2, cd.get("title", ""))
    else:
        half = len(charts) // 2 + len(charts) % 2
        positions = []
        for i in range(min(len(charts), 4)):
            row = i // 2
            col = i % 2
            positions.append((0.8 + col * 6.2, 1.5 + row * 2.9, 5.6, 2.7))
        for i, cd in enumerate(charts[:4]):
            chart_type_key = cd.get("type", "bar")
            xl_type = CHART_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)
            data = CategoryChartData()
            data.categories = cd.get("categories", [])
            for s in cd.get("series", []):
                data.add_series(s.get("name", ""), s.get("values", []))
            pos = positions[i]
            add_chart(slide, xl_type, data, pos[0], pos[1], pos[2], pos[3], cd.get("title", ""))

    if slide_data.get("notes"):
        add_textbox(slide, 0.8, 6.8, 11.7, 0.5, slide_data["notes"],
                    font_size=12, color=GRAY)


def build_summary_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BLUE_DARK)
    add_accent_bar(slide, 0, 3.2, 13.333, 0.08, BLUE_LIGHT)
    add_textbox(slide, 0.8, 0.8, 11.7, 0.8, slide_data.get("title", "Key Takeaways"),
                font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_accent_bar(slide, 5.6, 1.5, 2.0, 0.04, BLUE_LIGHT)

    takeaways = slide_data.get("takeaways", [])
    if takeaways:
        add_bullet_textbox(slide, 1.5, 1.8, 10.3, 3.5, takeaways,
                           font_size=18, color=WHITE)

    if slide_data.get("next_steps"):
        add_textbox(slide, 1.5, 5.2, 10.3, 0.5, "Next Steps",
                    font_size=16, bold=True, color=BLUE_LIGHT)
        add_bullet_textbox(slide, 1.5, 5.6, 10.3, 1.5, slide_data["next_steps"],
                           font_size=14, color=BLUE_LIGHT)


SLIDE_BUILDERS = {
    "title": build_title_slide,
    "content": build_content_slide,
    "chart": build_chart_slide,
    "summary": build_summary_slide,
}


def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_pptx.py <spec.json> [output_path]")
        sys.exit(1)

    spec_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "output/presentation.pptx"

    with open(spec_path, "r", encoding="utf-8") as f:
        spec = json.load(f)

    prs = create_presentation()
    slides = spec.get("slides", [])

    if len(slides) > 8:
        print(f"WARNING: {len(slides)} slides specified. Trimming to 8 max.")
        slides = slides[:8]

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type, build_content_slide)
        builder(prs, slide_data)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    print(f"OK: Saved {len(slides)} slides to {output_path}")


if __name__ == "__main__":
    main()
